"""Render an AWS architecture diagram into a .pptx using python-pptx.

Input is a YAML or JSON file. See SKILL.md for the schema. The slide layout is
16:9 (13.33" x 7.5"), with an AWS dark navy theme by default.

Usage:
    python render.py <input.yaml|json> [-o output.pptx]
"""
from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).parent))
from aws_catalog import Catalog  # noqa: E402

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# AWS official palette (Squid Ink + Smile)
AWS_NAVY = RGBColor(0x23, 0x2F, 0x3E)
AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)
AWS_LIGHT_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
AWS_MID_GRAY = RGBColor(0xA5, 0xA5, 0xA5)
AWS_DARK_GRAY = RGBColor(0x54, 0x5B, 0x64)
AWS_BLUE = RGBColor(0x53, 0x9F, 0xE5)
AWS_GREEN = RGBColor(0x28, 0xAD, 0x32)
AWS_RED = RGBColor(0xEB, 0x6F, 0x6F)

THEMES = {
    "dark": {
        "bg": AWS_NAVY,
        "fg": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_fg": AWS_LIGHT_GRAY,
        "node_fg": RGBColor(0xFF, 0xFF, 0xFF),
        "edge": AWS_LIGHT_GRAY,
        "group_line": AWS_ORANGE,
        "group_fg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "light": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "fg": AWS_NAVY,
        "subtitle_fg": AWS_DARK_GRAY,
        "node_fg": AWS_NAVY,
        "edge": AWS_DARK_GRAY,
        "group_line": AWS_ORANGE,
        "group_fg": AWS_NAVY,
    },
}

GROUP_BORDER_STYLE = {
    # Official AWS group container conventions (visual approximation)
    "vpc": {"color": (0x82, 0x82, 0xB1), "dash": "solid", "label": "VPC"},
    "aws-cloud": {"color": (0x23, 0x2F, 0x3E), "dash": "solid", "label": "AWS Cloud"},
    "aws-cloud-alt": {"color": (0x24, 0x8F, 0x33), "dash": "solid", "label": "AWS Cloud"},
    "region": {"color": (0x00, 0xA4, 0xA6), "dash": "dash", "label": "Region"},
    "az": {"color": (0x00, 0x73, 0xBB), "dash": "dash", "label": "Availability Zone"},
    "availability-zone": {"color": (0x00, 0x73, 0xBB), "dash": "dash", "label": "Availability Zone"},
    "public-subnet": {"color": (0x7A, 0xA1, 0x16), "dash": "solid", "label": "Public Subnet"},
    "private-subnet": {"color": (0x00, 0x73, 0xBB), "dash": "solid", "label": "Private Subnet"},
    "account": {"color": (0xCD, 0x21, 0x64), "dash": "dash", "label": "AWS Account"},
    "aws-account": {"color": (0xCD, 0x21, 0x64), "dash": "dash", "label": "AWS Account"},
    "auto-scaling-group": {"color": (0xED, 0x71, 0x00), "dash": "dash", "label": "Auto Scaling group"},
    "corporate-data-center": {"color": (0x7D, 0x8B, 0x99), "dash": "solid", "label": "Corporate Data Center"},
    "ec2-instance-contents": {"color": (0xED, 0x71, 0x00), "dash": "solid", "label": "EC2 Instance Contents"},
    "server-contents": {"color": (0x7D, 0x8B, 0x99), "dash": "solid", "label": "Server Contents"},
    "spot-fleet": {"color": (0xED, 0x71, 0x00), "dash": "dash", "label": "Spot Fleet"},
    # Generic / custom containers (caller can override colour & dash via custom_* fields)
    "custom": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": ""},
    "default": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": "Group"},
    # Pre-made colour palettes inspired by community AWS reference diagrams
    "operating-env": {"color": (0x88, 0x8C, 0x90), "dash": "solid", "label": "Operating ENV"},
    "compute-env": {"color": (0x88, 0x8C, 0x90), "dash": "solid", "label": "Compute ENV"},
    "next-gen-firewall": {"color": (0x6B, 0x70, 0x76), "dash": "solid", "label": "Next Generation Firewall"},
    "data-source": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": "Data Source"},
    "business-intelligence": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": "Business Intelligence"},
    "ai-bigdata": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": "AI/Big Data"},
    "report-generation": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": "Report Generation"},
    "workflow": {"color": (0xE7, 0x44, 0x7C), "dash": "solid", "label": ""},
    "workspaces": {"color": (0x6E, 0x73, 0x78), "dash": "dash", "label": "Workspaces"},
    "office": {"color": (0xA5, 0xA5, 0xA5), "dash": "dash", "label": "Office"},
    "highlight": {"color": (0xED, 0x71, 0x00), "dash": "solid", "label": ""},
}


@dataclass
class Node:
    id: str
    icon: str
    label: str = ""
    x: float | None = None  # inches
    y: float | None = None
    size: float = 0.85   # icon side in inches
    group: str | None = None
    sublabel: str = ""


@dataclass
class Group:
    id: str
    label: str = ""
    kind: str = "default"  # vpc / private-subnet / public-subnet / aws-cloud / region / az / account / ...
    parent: str | None = None
    nodes: list[str] = field(default_factory=list)
    children: list[str] = field(default_factory=list)
    pad: float = 0.45    # inches inside the border around contents
    direction: str = "auto"  # auto | row | column
    # Optional manual placement: any of these locks the group's bounding box
    # and skips auto-layout for it. Children inside still auto-flow within.
    x: float | None = None
    y: float | None = None
    w: float | None = None
    h: float | None = None
    # Custom styling overrides — let the caller paint a "Next Generation Firewall"
    # or "Data Source" container without inventing a new kind.
    border_color: str | None = None  # hex without "#" — e.g. "FF9900"
    border_style: str | None = None  # solid | dash | dot
    fill: str | None = None          # optional hex fill (inside the rectangle)
    fill_alpha: float = 0.0          # 0..1 opacity of the fill (default no fill)
    label_color: str | None = None
    label_size: int | None = None
    label_bold: bool = True
    label_italic: bool = False
    show_icon: bool = True           # show official group icon at top-left
    spans: list[str] = field(default_factory=list)  # ASG-style cross-subnet overlays


@dataclass
class Edge:
    source: str
    target: str
    label: str = ""
    style: str = "solid"  # solid | dash | dot
    direction: str = "->"  # -> | <-> | -
    color: str | None = None
    width: float = 1.5  # line thickness in points
    # Anchor side on each endpoint: "auto" (closest face) | "top" | "bottom" | "left" | "right"
    from_anchor: str = "auto"
    to_anchor: str = "auto"
    # Routing: "straight" (default) or "orthogonal" (auto L/Z) or "waypoints"
    routing: str = "straight"
    waypoints: list[tuple[float, float]] = field(default_factory=list)  # explicit bend points (inches)
    # Label appearance overrides
    label_color: str | None = None     # hex without "#"
    label_italic: bool = False
    label_bold: bool = False
    label_size: int = 9
    label_position: float = 0.5  # 0..1 along the edge (0=source, 1=target)
    label_offset: tuple[float, float] = (0.0, 0.0)  # nudge from natural position (inches)
    label_bg: str | None = None  # "transparent" | hex; default uses theme bg
    label_pad: float = 0.05


@dataclass
class Table:
    """Inline data table (used for route tables, IP plans, etc.)."""
    id: str
    x: float
    y: float
    w: float
    h: float
    columns: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    title: str = ""
    header_bg: str | None = None  # hex
    header_fg: str | None = None
    body_bg: str | None = None
    body_fg: str | None = None
    border_color: str | None = None
    font_size: int = 9


@dataclass
class Diagram:
    title: str = ""
    subtitle: str = ""
    description: str = ""
    theme: str = "dark"
    nodes: dict[str, Node] = field(default_factory=dict)
    groups: dict[str, Group] = field(default_factory=dict)
    edges: list[Edge] = field(default_factory=list)
    tables: list[Table] = field(default_factory=list)
    layout: str = "auto"  # auto | manual
    bg: str | None = None   # override theme bg color


# --------------------------------------------------------------------------- #
# Loading

def load_input(path: Path) -> list[Diagram]:
    raw = path.read_text(encoding="utf-8")
    data = yaml.safe_load(raw) if path.suffix.lower() in {".yaml", ".yml"} else json.loads(raw)

    if isinstance(data, dict) and "diagrams" in data:
        return [_parse_diagram(d) for d in data["diagrams"]]
    return [_parse_diagram(data)]


def _parse_diagram(d: dict) -> Diagram:
    diag = Diagram(
        title=d.get("title", ""),
        subtitle=d.get("subtitle", ""),
        description=d.get("description", ""),
        theme=d.get("theme", "dark"),
        layout=d.get("layout", "auto"),
        bg=d.get("bg"),
    )

    for g in d.get("groups", []) or []:
        gid = g["id"]
        diag.groups[gid] = Group(
            id=gid,
            label=g.get("label", g.get("kind", gid)),
            kind=g.get("kind", "default"),
            parent=g.get("parent"),
            pad=float(g.get("pad", 0.45)),
            direction=g.get("direction", "auto"),
            x=g.get("x"),
            y=g.get("y"),
            w=g.get("w"),
            h=g.get("h"),
            border_color=g.get("border_color"),
            border_style=g.get("border_style"),
            fill=g.get("fill"),
            fill_alpha=float(g.get("fill_alpha", 0.0)),
            label_color=g.get("label_color"),
            label_size=g.get("label_size"),
            label_bold=bool(g.get("label_bold", True)),
            label_italic=bool(g.get("label_italic", False)),
            show_icon=bool(g.get("show_icon", True)),
            spans=list(g.get("spans", []) or []),
        )
    # build child relationships
    for g in diag.groups.values():
        if g.parent and g.parent in diag.groups:
            diag.groups[g.parent].children.append(g.id)

    for n in d.get("nodes", []) or []:
        nid = n["id"]
        node = Node(
            id=nid,
            icon=n["icon"],
            label=n.get("label", nid.replace("_", " ").replace("-", " ")),
            sublabel=n.get("sublabel", ""),
            x=n.get("x"),
            y=n.get("y"),
            size=float(n.get("size", 0.85)),
            group=n.get("group"),
        )
        diag.nodes[nid] = node
        if node.group and node.group in diag.groups:
            diag.groups[node.group].nodes.append(nid)

    for e in d.get("edges", []) or []:
        wp = []
        for pt in (e.get("waypoints") or []):
            if isinstance(pt, dict):
                wp.append((float(pt["x"]), float(pt["y"])))
            elif isinstance(pt, (list, tuple)) and len(pt) >= 2:
                wp.append((float(pt[0]), float(pt[1])))
        offset = e.get("label_offset") or (0.0, 0.0)
        if isinstance(offset, dict):
            offset = (float(offset.get("x", 0.0)), float(offset.get("y", 0.0)))
        else:
            offset = (float(offset[0]), float(offset[1])) if len(offset) >= 2 else (0.0, 0.0)
        diag.edges.append(Edge(
            source=e["from"],
            target=e["to"],
            label=e.get("label", ""),
            style=e.get("style", "solid"),
            direction=e.get("direction", "->"),
            color=e.get("color"),
            width=float(e.get("width", 1.5)),
            from_anchor=e.get("from_anchor", "auto"),
            to_anchor=e.get("to_anchor", "auto"),
            routing=e.get("routing", "waypoints" if wp else "straight"),
            waypoints=wp,
            label_color=e.get("label_color"),
            label_italic=bool(e.get("label_italic", False)),
            label_bold=bool(e.get("label_bold", False)),
            label_size=int(e.get("label_size", 9)),
            label_position=float(e.get("label_position", 0.5)),
            label_offset=offset,
            label_bg=e.get("label_bg"),
            label_pad=float(e.get("label_pad", 0.05)),
        ))

    for t in d.get("tables", []) or []:
        diag.tables.append(Table(
            id=t.get("id", f"tbl-{len(diag.tables)}"),
            x=float(t["x"]),
            y=float(t["y"]),
            w=float(t["w"]),
            h=float(t["h"]),
            columns=list(t.get("columns", []) or []),
            rows=[list(r) for r in (t.get("rows", []) or [])],
            title=t.get("title", ""),
            header_bg=t.get("header_bg"),
            header_fg=t.get("header_fg"),
            body_bg=t.get("body_bg"),
            body_fg=t.get("body_fg"),
            border_color=t.get("border_color"),
            font_size=int(t.get("font_size", 9)),
        ))

    return diag


# --------------------------------------------------------------------------- #
# Auto-layout

def _content_weight(diag: Diagram, kind: str, ident: str) -> int:
    """Heuristic 'how much space this unit should claim' = leaf node count."""
    if kind == "node":
        return 1
    g = diag.groups[ident]
    own_nodes = sum(1 for nid in g.nodes if diag.nodes[nid].x is None)
    children = sum(_content_weight(diag, "group", c) for c in g.children)
    return max(1, own_nodes + children)


def _is_manually_placed(g: Group) -> bool:
    """A group is manually placed if all four of x/y/w/h are set."""
    return all(v is not None for v in (g.x, g.y, g.w, g.h))


def _apply_manual_groups(diag: Diagram) -> None:
    """Place manually-positioned groups (any depth) and recurse into their contents.

    After this runs, `_x/_y/_w/_h` are set on those groups, and their children
    (nodes + nested groups) are positioned within their inner area.
    """
    for g in diag.groups.values():
        if not _is_manually_placed(g):
            continue
        # Compute the center to feed into _layout_group, which fills _x/_y/_w/_h
        # and recurses into the children. Children that are themselves manually
        # placed will be re-placed by their own pass below — that's fine because
        # _layout_group respects already-set node.x positions.
        _layout_group(diag, g.id,
                      cx=g.x + g.w / 2,  # type: ignore[operator]
                      cy=g.y + g.h / 2,  # type: ignore[operator]
                      w=g.w, h=g.h,  # type: ignore[arg-type]
                      depth=_group_depth(diag, g))


def _apply_span_groups(diag: Diagram) -> None:
    """Position groups that 'span' (overlay) other already-placed groups.

    Used for Auto Scaling Group overlays that visually wrap multiple subnets.
    A spanning group draws on top of its children but doesn't re-layout them.
    """
    for g in diag.groups.values():
        if not g.spans:
            continue
        if hasattr(g, "_x"):
            continue
        xs: list[float] = []
        ys: list[float] = []
        xe: list[float] = []
        ye: list[float] = []
        for sid in g.spans:
            sg = diag.groups.get(sid)
            if not sg or not hasattr(sg, "_x"):
                continue
            xs.append(sg._x)  # type: ignore[attr-defined]
            ys.append(sg._y)  # type: ignore[attr-defined]
            xe.append(sg._x + sg._w)  # type: ignore[attr-defined]
            ye.append(sg._y + sg._h)  # type: ignore[attr-defined]
        # also include nodes the spanning group lists directly
        for nid in g.nodes:
            n = diag.nodes.get(nid)
            if n and n.x is not None and n.y is not None:
                xs.append(n.x - 0.1)
                ys.append(n.y - 0.1)
                xe.append(n.x + n.size + 0.1)
                ye.append(n.y + n.size + 0.4)
        if not xs:
            continue
        pad = max(g.pad - 0.1, 0.10)
        g._x = min(xs) - pad  # type: ignore[attr-defined]
        g._y = min(ys) - pad - 0.18  # type: ignore[attr-defined]
        g._w = max(xe) - min(xs) + 2 * pad  # type: ignore[attr-defined]
        g._h = max(ye) - min(ys) + 2 * pad + 0.18  # type: ignore[attr-defined]


def _group_depth(diag: Diagram, g: Group) -> int:
    d = 0
    cur = g
    while cur.parent and cur.parent in diag.groups:
        d += 1
        cur = diag.groups[cur.parent]
    return d


def auto_layout(diag: Diagram, area: tuple[float, float, float, float]) -> None:
    """Lay out nodes that don't have explicit x/y. Pack groups + their members.

    area = (x, y, w, h) in inches reserved for the diagram body.
    """
    ax, ay, aw, ah = area

    # First, fix any group with manual x/y/w/h (recurses into its children).
    _apply_manual_groups(diag)

    top_groups = [g for g in diag.groups.values()
                  if not g.parent and not _is_manually_placed(g)]
    ungrouped_nodes = [n for n in diag.nodes.values() if not n.group and n.x is None]

    units: list[tuple[str, str]] = [("group", g.id) for g in top_groups] + [
        ("node", n.id) for n in ungrouped_nodes
    ]
    if not units:
        return

    _layout_units(diag, units, ax, ay, aw, ah, depth=0, direction="auto")


def _layout_units(diag: Diagram, units: list[tuple[str, str]],
                  x: float, y: float, w: float, h: float, depth: int,
                  direction: str = "auto") -> None:
    """Place a list of units inside (x, y, w, h). Used for both root and inside groups."""
    import math
    if not units:
        return

    n = len(units)
    weights = [max(1, _content_weight(diag, k, i)) for k, i in units]
    total_w = sum(weights)

    aspect = w / max(h, 0.01)
    if direction == "row":
        cols, rows = n, 1
    elif direction == "column":
        cols, rows = 1, n
    elif n == 1:
        cols, rows = 1, 1
    elif n <= 4 and aspect >= 1.3:
        cols, rows = n, 1
    elif n <= 4 and aspect <= 0.77:
        cols, rows = 1, n
    else:
        cols = max(1, min(n, math.ceil(math.sqrt(n * aspect))))
        rows = math.ceil(n / cols)

    # In single-row/single-col mode, allocate proportional widths/heights.
    if rows == 1:
        cx = x
        for (kind, ident), wt in zip(units, weights):
            cell_w = w * (wt / total_w)
            _place_unit(diag, kind, ident,
                        cx + cell_w / 2, y + h / 2,
                        cell_w * 0.95, h * 0.95, depth)
            cx += cell_w
        return
    if cols == 1:
        cy = y
        for (kind, ident), wt in zip(units, weights):
            cell_h = h * (wt / total_w)
            _place_unit(diag, kind, ident,
                        x + w / 2, cy + cell_h / 2,
                        w * 0.95, cell_h * 0.95, depth)
            cy += cell_h
        return

    # Generic grid (uniform cells; weights ignored for simplicity).
    cell_w = w / cols
    cell_h = h / rows
    for idx, (kind, ident) in enumerate(units):
        c = idx % cols
        r = idx // cols
        ccx = x + c * cell_w + cell_w / 2
        ccy = y + r * cell_h + cell_h / 2
        _place_unit(diag, kind, ident, ccx, ccy,
                    cell_w * 0.95, cell_h * 0.92, depth)


def _place_unit(diag: Diagram, kind: str, ident: str,
                cx: float, cy: float, w: float, h: float, depth: int) -> None:
    if kind == "group":
        _layout_group(diag, ident, cx, cy,
                      max(w, 1.0), max(h, 0.9), depth + 1)
    else:
        node = diag.nodes[ident]
        # Shrink node icon if cell is too small so it doesn't overflow its group.
        # Reserve ~0.32" below for the node label.
        max_side = max(0.42, min(node.size, w - 0.05, h - 0.32))
        node.size = max_side
        node.x = cx - node.size / 2
        node.y = cy - node.size / 2 - 0.12


def _layout_group(diag: Diagram, gid: str,
                  cx: float, cy: float, w: float, h: float, depth: int) -> None:
    """Place a group at (cx, cy) with width w, height h, recursing into children."""
    group = diag.groups[gid]
    group._x = cx - w / 2  # type: ignore[attr-defined]
    group._y = cy - h / 2  # type: ignore[attr-defined]
    group._w = w  # type: ignore[attr-defined]
    group._h = h  # type: ignore[attr-defined]

    # Padding/label shrink with depth so deep nesting doesn't eat all the space.
    pad = max(group.pad - depth * 0.07, 0.16)
    label_h = max(0.40 - depth * 0.04, 0.26)
    inner_x = group._x + pad  # type: ignore[attr-defined]
    inner_y = group._y + pad + label_h  # type: ignore[attr-defined]
    inner_w = max(w - 2 * pad, 0.4)
    inner_h = max(h - 2 * pad - label_h, 0.4)

    # Stash chosen pad so the renderer can use it later if needed.
    group._pad = pad  # type: ignore[attr-defined]
    group._label_h = label_h  # type: ignore[attr-defined]

    units: list[tuple[str, str]] = [
        ("group", c) for c in group.children
        if not _is_manually_placed(diag.groups[c])
    ] + [
        ("node", nid) for nid in group.nodes
        if diag.nodes[nid].x is None
    ]
    if not units:
        return

    _layout_units(diag, units, inner_x, inner_y, inner_w, inner_h, depth,
                  direction=group.direction)


# --------------------------------------------------------------------------- #
# Rendering

def render(diag: Diagram, prs: Presentation, catalog: Catalog) -> None:
    theme_name = diag.theme if diag.theme in THEMES else "dark"
    theme = THEMES[theme_name]

    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(diag.bg) if diag.bg else theme["bg"]
    bg.shadow.inherit = False

    # Title bar
    title_h = Inches(0.7) if diag.title else Inches(0)
    if diag.title:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_W - Inches(1), Inches(0.6))
        tf = tb.text_frame
        tf.margin_top = tf.margin_bottom = Inches(0.05)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = diag.title
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = theme["fg"]
        run.font.name = "Amazon Ember"
        if diag.subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = diag.subtitle
            r2.font.size = Pt(13)
            r2.font.color.rgb = theme["subtitle_fg"]
            r2.font.name = "Amazon Ember"

    # Diagram body area (inches)
    body_top = 1.05 if diag.title else 0.5
    body = (0.6, body_top, 13.333 - 1.2, 7.5 - body_top - 0.4)

    # Auto-place
    if diag.layout != "manual":
        auto_layout(diag, body)

    # Resolve spanning groups (ASG overlays) once subnets/nodes have positions
    _apply_span_groups(diag)

    # Draw groups (largest first so children render on top)
    drawn_groups: list[Group] = []
    for g in diag.groups.values():
        if hasattr(g, "_x"):
            drawn_groups.append(g)
    # also handle groups whose position derives from member nodes (manual layout)
    for g in diag.groups.values():
        if not hasattr(g, "_x"):
            _derive_group_bounds(diag, g)
            if hasattr(g, "_x"):
                drawn_groups.append(g)

    drawn_groups.sort(key=lambda g: -(getattr(g, "_w", 0) * getattr(g, "_h", 0)))
    for g in drawn_groups:
        _draw_group(slide, g, theme, catalog)

    # Draw nodes
    for n in diag.nodes.values():
        if n.x is None or n.y is None:
            continue
        _draw_node(slide, n, theme, catalog)

    # Draw edges
    for e in diag.edges:
        _draw_edge(slide, e, diag, theme)

    # Draw inline tables (e.g. route tables) on top of everything
    for t in diag.tables:
        _draw_table(slide, t, theme)

    # Footer description
    if diag.description:
        fb = slide.shapes.add_textbox(Inches(0.5), SLIDE_H - Inches(0.45),
                                      SLIDE_W - Inches(1), Inches(0.4))
        ftf = fb.text_frame
        ftf.margin_top = Inches(0)
        p = ftf.paragraphs[0]
        run = p.add_run()
        run.text = diag.description
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = theme["subtitle_fg"]


def _derive_group_bounds(diag: Diagram, g: Group) -> None:
    """For manual layouts, compute group bounds from member node positions."""
    xs: list[float] = []
    ys: list[float] = []
    xe: list[float] = []
    ye: list[float] = []
    for nid in g.nodes:
        n = diag.nodes.get(nid)
        if n and n.x is not None and n.y is not None:
            xs.append(n.x)
            ys.append(n.y)
            xe.append(n.x + n.size)
            ye.append(n.y + n.size + 0.3)
    for cid in g.children:
        c = diag.groups.get(cid)
        if c and hasattr(c, "_x"):
            xs.append(c._x)  # type: ignore[attr-defined]
            ys.append(c._y)  # type: ignore[attr-defined]
            xe.append(c._x + c._w)  # type: ignore[attr-defined]
            ye.append(c._y + c._h)  # type: ignore[attr-defined]
    if not xs:
        return
    pad = g.pad
    g._x = min(xs) - pad  # type: ignore[attr-defined]
    g._y = min(ys) - pad - 0.42  # type: ignore[attr-defined]
    g._w = max(xe) - min(xs) + 2 * pad  # type: ignore[attr-defined]
    g._h = max(ye) - min(ys) + 2 * pad + 0.42  # type: ignore[attr-defined]


def _draw_group(slide, g: Group, theme: dict, catalog: Catalog) -> None:
    style = GROUP_BORDER_STYLE.get(g.kind, GROUP_BORDER_STYLE["default"])
    x = Inches(g._x)  # type: ignore[attr-defined]
    y = Inches(g._y)  # type: ignore[attr-defined]
    w = Inches(g._w)  # type: ignore[attr-defined]
    h = Inches(g._h)  # type: ignore[attr-defined]

    # Resolve effective border colour & dash. Custom overrides win.
    if g.border_color:
        line_rgb = _hex_to_rgb(g.border_color)
    else:
        line_rgb = RGBColor(*style["color"])
    dash_kind = (g.border_style or style["dash"]).lower()

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if g.fill and g.fill_alpha > 0:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(g.fill)
        # python-pptx doesn't expose alpha directly; emulate via transparency XML
        try:
            from lxml import etree
            spPr = shape.fill._xPr.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
            )
            if spPr is not None:
                srgb = spPr.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                )
                if srgb is not None:
                    alpha_pct = int(max(0.0, min(1.0, g.fill_alpha)) * 100000)
                    alpha = etree.SubElement(
                        srgb,
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
                    )
                    alpha.set("val", str(alpha_pct))
        except Exception:
            pass
    else:
        shape.fill.background()

    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.25)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        if dash_kind == "dash":
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif dash_kind in ("dot", "dotted"):
            shape.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        elif dash_kind in ("dash-dot", "dashdot"):
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT
        elif dash_kind == "long-dash":
            shape.line.dash_style = MSO_LINE_DASH_STYLE.LONG_DASH
    except Exception:
        pass
    shape.shadow.inherit = False
    # Suppress default text frame
    if shape.has_text_frame:
        shape.text_frame.text = ""

    # Try to drop the official group icon on top-left of border
    icon_path: Path | None = None
    icon_keys = {
        "vpc": "grp-virtual-private-cloud-vpc",
        "aws-cloud": "grp-aws-cloud",
        "aws-cloud-alt": "grp-aws-cloud",
        "region": "grp-region",
        "public-subnet": "grp-public-subnet",
        "private-subnet": "grp-private-subnet",
        "account": "grp-aws-account",
        "aws-account": "grp-aws-account",
        "auto-scaling-group": "grp-auto-scaling-group",
        "corporate-data-center": "grp-corporate-data-center",
        "ec2-instance-contents": "grp-ec2-instance-contents",
        "server-contents": "grp-server-contents",
        "spot-fleet": "grp-spot-fleet",
    }
    if g.show_icon and g.kind in icon_keys and icon_keys[g.kind] in catalog.icons:
        icon_path = catalog.root / catalog.icons[icon_keys[g.kind]]["path"]

    width_in = g._w  # type: ignore[attr-defined]
    show_icon = bool(icon_path) and width_in >= 1.6 and g.show_icon
    label_x = x + Inches(0.18)
    if show_icon and icon_path:
        icon_w = Inches(0.30)
        slide.shapes.add_picture(str(icon_path), x + Inches(0.14), y + Inches(0.10),
                                 width=icon_w, height=icon_w)
        label_x = x + Inches(0.50)

    label = g.label or style["label"]
    if not label:
        return
    label_w_in = max(width_in - (0.55 if show_icon else 0.30), 0.5)
    tb = slide.shapes.add_textbox(label_x, y + Inches(0.06), Inches(label_w_in), Inches(0.32))
    tf = tb.text_frame
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.margin_left = tf.margin_right = Inches(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    # shrink label font for narrow groups so it doesn't overflow visibly
    if g.label_size:
        run.font.size = Pt(g.label_size)
    else:
        run.font.size = Pt(11) if width_in >= 2.0 else Pt(9)
    run.font.bold = g.label_bold
    run.font.italic = g.label_italic
    run.font.color.rgb = _hex_to_rgb(g.label_color) if g.label_color else line_rgb
    run.font.name = "Amazon Ember"


def _draw_table(slide, t: Table, theme: dict) -> None:
    """Draw a small data table — used for inline route tables, IP plans, etc."""
    rows_total = len(t.rows) + (1 if t.columns else 0) + (1 if t.title else 0)
    cols_total = max(len(t.columns), max((len(r) for r in t.rows), default=1))
    if rows_total == 0 or cols_total == 0:
        return

    tbl_shape = slide.shapes.add_table(
        rows_total, cols_total,
        Inches(t.x), Inches(t.y), Inches(t.w), Inches(t.h),
    )
    table = tbl_shape.table
    border_rgb = _hex_to_rgb(t.border_color) if t.border_color else theme["edge"]
    header_bg = _hex_to_rgb(t.header_bg) if t.header_bg else _hex_to_rgb("3B4859")
    header_fg = _hex_to_rgb(t.header_fg) if t.header_fg else theme["fg"]
    body_bg = _hex_to_rgb(t.body_bg) if t.body_bg else _hex_to_rgb("2A3645")
    body_fg = _hex_to_rgb(t.body_fg) if t.body_fg else theme["fg"]

    row_idx = 0
    if t.title:
        # Merge first row for title
        row = table.rows[row_idx]
        cell0 = row.cells[0]
        for c in range(1, cols_total):
            cell0.merge(row.cells[c])
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = header_bg
        tf = cell0.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = 2
        run = p.add_run()
        run.text = t.title
        run.font.bold = True
        run.font.size = Pt(t.font_size + 1)
        run.font.color.rgb = header_fg
        run.font.name = "Amazon Ember"
        row_idx += 1

    if t.columns:
        row = table.rows[row_idx]
        for c, name in enumerate(t.columns):
            cell = row.cells[c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = 2
            run = p.add_run()
            run.text = str(name)
            run.font.bold = True
            run.font.size = Pt(t.font_size)
            run.font.color.rgb = header_fg
            run.font.name = "Amazon Ember"
        row_idx += 1

    for r, row_data in enumerate(t.rows):
        row = table.rows[row_idx + r]
        for c in range(cols_total):
            cell = row.cells[c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = body_bg
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = 2
            run = p.add_run()
            run.text = str(row_data[c]) if c < len(row_data) else ""
            run.font.size = Pt(t.font_size)
            run.font.color.rgb = body_fg
            run.font.name = "Amazon Ember"


def _draw_node(slide, n: Node, theme: dict, catalog: Catalog) -> None:
    try:
        path = catalog.resolve(n.icon)
    except KeyError:
        suggestions = catalog.suggest(n.icon, limit=5)
        msg = f"Unknown icon '{n.icon}'."
        if suggestions:
            msg += f" Try: {', '.join(suggestions)}"
        raise SystemExit(msg)

    x = Inches(n.x)  # type: ignore[arg-type]
    y = Inches(n.y)  # type: ignore[arg-type]
    side = Inches(n.size)
    slide.shapes.add_picture(str(path), x, y, width=side, height=side)

    # Label below
    label_w = Inches(max(n.size + 0.6, 1.4))
    label_x = x - (label_w - side) / 2
    tb = slide.shapes.add_textbox(label_x, y + side + Inches(0.04), label_w, Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.margin_left = tf.margin_right = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = n.label
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = theme["node_fg"]
    run.font.name = "Amazon Ember"

    if n.sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = 2
        r2 = p2.add_run()
        r2.text = n.sublabel
        r2.font.size = Pt(8)
        r2.font.color.rgb = theme["subtitle_fg"]
        r2.font.name = "Amazon Ember"


def _node_center(n: Node) -> tuple[float, float]:
    return (n.x + n.size / 2, n.y + n.size / 2)  # type: ignore[operator]


def _node_anchor(src: Node, dst: Node, side: str = "auto") -> tuple[float, float]:
    """Pick an anchor on src's bounding box.

    side: "auto" picks the face nearest to dst; "top"/"bottom"/"left"/"right"
    forces a specific face.
    """
    sx, sy = _node_center(src)
    dx, dy = _node_center(dst)
    half = src.size / 2
    if side == "top":
        return (sx, sy - half)
    if side == "bottom":
        return (sx, sy + half)
    if side == "left":
        return (sx - half, sy)
    if side == "right":
        return (sx + half, sy)
    if abs(dx - sx) >= abs(dy - sy):
        return (sx + (half if dx >= sx else -half), sy)
    return (sx, sy + (half if dy >= sy else -half))


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_dash(line, style: str) -> None:
    if style not in ("dash", "dot", "dotted", "dash-dot", "dashdot", "long-dash", "solid"):
        return
    if style == "solid":
        return
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        if style == "dash":
            line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif style in ("dot", "dotted"):
            line.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        elif style in ("dash-dot", "dashdot"):
            line.line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT
        elif style == "long-dash":
            line.line.dash_style = MSO_LINE_DASH_STYLE.LONG_DASH
    except Exception:
        pass


def _build_edge_path(e: Edge, src: Node, dst: Node) -> list[tuple[float, float]]:
    """Return the polyline of (x, y) inch points for the connector.

    Length is at least 2. Waypoints win over orthogonal which wins over straight.
    """
    sx, sy = _node_anchor(src, dst, e.from_anchor)
    dx, dy = _node_anchor(dst, src, e.to_anchor)
    if e.routing == "waypoints" and e.waypoints:
        # Anchor toward the first waypoint, end-anchor toward the last waypoint
        first = e.waypoints[0]
        last = e.waypoints[-1]
        sx, sy = _node_anchor_toward(src, first[0], first[1], e.from_anchor)
        dx, dy = _node_anchor_toward(dst, last[0], last[1], e.to_anchor)
        return [(sx, sy), *e.waypoints, (dx, dy)]
    if e.routing == "orthogonal":
        # 2-bend L/Z routing depending on dominant axis
        # Decide first leg based on the chosen anchor face
        face = e.from_anchor if e.from_anchor != "auto" else _dominant_face(src, dst)
        if face in ("left", "right"):
            # horizontal first
            mx = (sx + dx) / 2
            return [(sx, sy), (mx, sy), (mx, dy), (dx, dy)]
        else:
            # vertical first
            my = (sy + dy) / 2
            return [(sx, sy), (sx, my), (dx, my), (dx, dy)]
    return [(sx, sy), (dx, dy)]


def _node_anchor_toward(n: Node, tx: float, ty: float, side: str) -> tuple[float, float]:
    sx = n.x + n.size / 2  # type: ignore[operator]
    sy = n.y + n.size / 2  # type: ignore[operator]
    half = n.size / 2
    if side == "top":
        return (sx, sy - half)
    if side == "bottom":
        return (sx, sy + half)
    if side == "left":
        return (sx - half, sy)
    if side == "right":
        return (sx + half, sy)
    if abs(tx - sx) >= abs(ty - sy):
        return (sx + (half if tx >= sx else -half), sy)
    return (sx, sy + (half if ty >= sy else -half))


def _dominant_face(src: Node, dst: Node) -> str:
    sx, sy = _node_center(src)
    dx, dy = _node_center(dst)
    if abs(dx - sx) >= abs(dy - sy):
        return "right" if dx >= sx else "left"
    return "bottom" if dy >= sy else "top"


def _draw_polyline(slide, points: list[tuple[float, float]], color: RGBColor,
                   width_pt: float, style: str,
                   head: bool, tail: bool):
    """Draw a polyline as a series of straight connectors. Returns the last connector.
    Arrowheads only on the first/last segment so the path appears continuous."""
    n = len(points)
    last = None
    for i in range(n - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color
        line.line.width = Pt(width_pt)
        _apply_dash(line, style)
        # Suppress fill / shadow
        try:
            line.shadow.inherit = False
        except Exception:
            pass
        seg_head = head and (i == n - 2)
        seg_tail = tail and (i == 0)
        if seg_head or seg_tail:
            _set_arrow(line, head=seg_head, tail=seg_tail)
        else:
            _set_arrow(line, head=False, tail=False)
        last = line
    return last


def _polyline_midpoint(points: list[tuple[float, float]], frac: float) -> tuple[float, float]:
    """Point a fractional distance along the polyline (frac in 0..1)."""
    if len(points) < 2:
        return points[0] if points else (0.0, 0.0)
    seg_lens = []
    for i in range(len(points) - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        seg_lens.append(((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5)
    total = sum(seg_lens) or 1.0
    target = frac * total
    acc = 0.0
    for i, l in enumerate(seg_lens):
        if acc + l >= target:
            t = (target - acc) / l if l > 0 else 0
            x1, y1 = points[i]
            x2, y2 = points[i + 1]
            return (x1 + (x2 - x1) * t, y1 + (y2 - y1) * t)
        acc += l
    return points[-1]


def _draw_edge(slide, e: Edge, diag: Diagram, theme: dict) -> None:
    src = diag.nodes.get(e.source)
    dst = diag.nodes.get(e.target)
    if not src or not dst or src.x is None or dst.x is None:
        return

    points = _build_edge_path(e, src, dst)
    color = _hex_to_rgb(e.color) if e.color else theme["edge"]
    head = e.direction in {"->", "<->"}
    tail = e.direction == "<->"
    _draw_polyline(slide, points, color, e.width, e.style, head=head, tail=tail)

    if e.label:
        lx, ly = _polyline_midpoint(points, max(0.0, min(1.0, e.label_position)))
        ox, oy = e.label_offset
        lx += ox
        ly += oy
        tw = max(0.8, min(3.6, 0.085 * len(e.label) + 0.5))
        tb = slide.shapes.add_textbox(Inches(lx - tw / 2), Inches(ly - 0.16),
                                      Inches(tw), Inches(0.32))
        tf = tb.text_frame
        tf.margin_top = tf.margin_bottom = Inches(e.label_pad)
        tf.margin_left = tf.margin_right = Inches(0.05)
        # Background: caller can request transparent
        if e.label_bg == "transparent" or e.label_bg == "none":
            tb.fill.background()
        elif e.label_bg:
            tb.fill.solid()
            tb.fill.fore_color.rgb = _hex_to_rgb(e.label_bg)
        else:
            tb.fill.solid()
            tb.fill.fore_color.rgb = theme["bg"]
        tb.line.fill.background()
        p = tf.paragraphs[0]
        p.alignment = 2
        run = p.add_run()
        run.text = e.label
        run.font.size = Pt(e.label_size)
        run.font.color.rgb = _hex_to_rgb(e.label_color) if e.label_color else theme["fg"]
        run.font.name = "Amazon Ember"
        if e.label_italic:
            run.font.italic = True
        if e.label_bold:
            run.font.bold = True


def _set_arrow(connector, head: bool, tail: bool) -> None:
    """Toggle arrowheads on a python-pptx connector via raw XML."""
    from lxml import etree
    spPr = connector.line._get_or_add_ln()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    # Remove any existing head/tail
    for tag in ("headEnd", "tailEnd"):
        for el in spPr.findall(f"a:{tag}", nsmap):
            spPr.remove(el)
    if tail:
        head_el = etree.SubElement(spPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd")
        head_el.set("type", "triangle")
        head_el.set("w", "med")
        head_el.set("len", "med")
    if head:
        tail_el = etree.SubElement(spPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
        tail_el.set("type", "triangle")
        tail_el.set("w", "med")
        tail_el.set("len", "med")


# --------------------------------------------------------------------------- #
# Entry point

def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("input", type=Path, help="YAML or JSON architecture spec")
    ap.add_argument("-o", "--output", type=Path, default=None,
                    help="Output .pptx path (default: <input>.pptx)")
    args = ap.parse_args()

    if not args.input.exists():
        print(f"Input not found: {args.input}", file=sys.stderr)
        return 1

    out = args.output or args.input.with_suffix(".pptx")

    diagrams = load_input(args.input)
    catalog = Catalog()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for diag in diagrams:
        render(diag, prs, catalog)

    prs.save(str(out))
    print(f"Wrote {out} ({len(diagrams)} slide(s))")
    return 0


if __name__ == "__main__":
    sys.exit(main())
